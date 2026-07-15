import sys
import os

def extract_pdf(path):
    import pypdf
    text = []
    try:
        reader = pypdf.PdfReader(path)
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text.append(t)
    except Exception as e:
        sys.stderr.write(f"PDF Error: {str(e)}\n")
    return "\n".join(text)

def extract_docx(path):
    import docx
    text = []
    try:
        doc = docx.Document(path)
        for para in doc.paragraphs:
            if para.text:
                text.append(para.text)
    except Exception as e:
        sys.stderr.write(f"DOCX Error: {str(e)}\n")
    return "\n".join(text)

def extract_pptx(path):
    from pptx import Presentation
    text = []
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
    except Exception as e:
        sys.stderr.write(f"PPTX Error: {str(e)}\n")
    return "\n".join(text)

def main():
    if len(sys.argv) < 2:
        sys.exit(1)
        
    path = sys.argv[1]
    if not os.path.exists(path):
        sys.exit(1)
        
    ext = os.path.splitext(path)[1].lower()
    if ext == '.pdf':
        text = extract_pdf(path)
    elif ext == '.docx':
        text = extract_docx(path)
    elif ext == '.pptx':
        text = extract_pptx(path)
    else:
        sys.exit(1)
        
    try:
        sys.stdout.reconfigure(encoding='utf-8')
    except AttributeError:
        pass
    sys.stdout.write(text)

if __name__ == '__main__':
    main()
